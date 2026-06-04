"""答辩PPT最终版 v2：图片用宽度约束，信息密集"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os, shutil

new_pptx = 'D:/bylw/code/答辩PPT/答辩PPT_v2.pptx'
fig = 'D:/bylw/code/lunwen/QLULatex/QLUThesisLatexTemplate-master/Thesis/static/figures'
mlb = 'D:/bylw/code/fangzhen/matlab'

prs = Presentation()
prs.slide_width = 12192000; prs.slide_height = 6858000
SW=33.87; SH=19.05

NAVY=RGBColor(0x1B,0x2A,0x4A); BLUE=RGBColor(0x2E,0x4A,0x7A)
LBLUE=RGBColor(0x3D,0x6C,0xB9); WH=RGBColor(0xFF,0xFF,0xFF)
OFF=RGBColor(0xF5,0xF7,0xFA); GR=RGBColor(0x6B,0x7B,0x8D)
DK=RGBColor(0x2C,0x3E,0x50); LGR=RGBColor(0xE8,0xEC,0xF1)

def ns():
    for l in prs.slide_layouts:
        if l.name=='Blank' or l.placeholders==0: return prs.slides.add_slide(l)
    return prs.slides.add_slide(prs.slide_layouts[0])

def R(sl,l,t,w,h,f,li=None):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Cm(l),Cm(t),Cm(w),Cm(h))
    s.fill.solid();s.fill.fore_color.rgb=f
    if li: s.line.color.rgb=li;s.line.width=Pt(1)
    else: s.line.fill.background()
    return s

def T(sl,l,t,w,h,txt,sz=12,b=False,c=DK,a=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Cm(l),Cm(t),Cm(w),Cm(h))
    tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.bold=b;p.font.color.rgb=c;p.alignment=a

def I(sl,path,l,t,w=None,h=None,max_h=None):
    """智能约束: 如果设了max_h且按w算出的高度超了max_h,就改用max_h"""
    if w and h: return sl.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w),height=Cm(h))
    elif w and max_h:
        from PIL import Image as PILImage
        im=PILImage.open(path); ratio=im.width/im.height
        need_h=w/ratio
        if need_h>max_h: return sl.shapes.add_picture(path,Cm(l),Cm(t),height=Cm(max_h))
        return sl.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w))
    elif w: return sl.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w))
    elif h: return sl.shapes.add_picture(path,Cm(l),Cm(t),height=Cm(h))
    else: return sl.shapes.add_picture(path,Cm(l),Cm(t))

def ST(sl,num,title,pg):
    R(sl,0,0,SW,0.08,LBLUE)
    T(sl,1,0.4,4,0.7,f'PART  {num:02d}',sz=13,c=GR)
    T(sl,1,1.1,25,0.8,title,sz=22,b=True,c=NAVY)
    R(sl,1,2.0,5,0.04,LBLUE)
    T(sl,30,17.8,3,0.5,pg,sz=9,c=GR,a=PP_ALIGN.RIGHT)

def FT(sl):
    R(sl,0,18.2,SW,0.85,OFF)
    T(sl,1,18.3,20,0.5,'齐鲁工业大学  机械工程学院',sz=8,c=GR)

# ═══ Slide 1: 封面 ═══
s=ns(); R(s,0,0,SW,SH,NAVY)
R(s,2,4,8,0.04,LBLUE)
T(s,2,4.4,30,1.5,'智能仓储移动机器人',sz=34,b=True,c=WH)
T(s,2,6.0,30,0.9,'结构设计与运动控制研究',sz=21,c=RGBColor(0x90,0xCA,0xF9))
T(s,2,7.5,30,0.7,'Design and Control of Intelligent Warehouse Mobile Robot',sz=12,c=RGBColor(0x78,0x90,0x9C))
R(s,2,8.8,8,0.04,LBLUE)
T(s,2,9.5,15,0.5,'答辩人：王梓煜',sz=12,c=RGBColor(0xBB,0xDE,0xFB))
T(s,2,10.2,15,0.5,'指导教师：曹树春',sz=12,c=RGBColor(0xBB,0xDE,0xFB))
T(s,2,10.9,15,0.5,'齐鲁工业大学  机械工程学院',sz=12,c=RGBColor(0xBB,0xDE,0xFB))
T(s,2,12.0,15,0.5,'2026年6月',sz=13,b=True,c=WH)
logo=os.path.join(fig,'MainLOGO-pretty.png')
if os.path.exists(logo):
    try: I(s,logo,20,9.5,w=12,max_h=3.5)
    except: pass

# ═══ Slide 2: 目录 ═══
s=ns(); R(s,0,0,SW,SH,WH); R(s,0,0,SW,0.08,LBLUE)
T(s,2,0.8,10,0.9,'目  录',sz=26,b=True,c=NAVY)
T(s,2,1.7,10,0.4,'CONTENTS',sz=11,c=GR)
R(s,2,2.2,3.5,0.04,LBLUE)
toc=[('01','研究背景与意义','研究现状、目标与方法'),('02','机械结构设计','总装布局、升降机构、强度校核'),('03','电气系统设计','主控、电机、传感器、电路'),('04','运动学建模与仿真','差速模型、四类轨迹验证'),('05','路径规划与循迹','Hybrid A* + PID红外循迹'),('06','运动基元与算法','5种曲率基元、算法流程'),('07','双环PID控制','级联结构、四项仿真验证')]
for i,(n,na,de) in enumerate(toc):
    y=3.0+i*2.1; R(s,2,y,30,1.7,OFF); R(s,2,y,0.12,1.7,LBLUE)
    T(s,3,y+0.15,2,0.9,n,sz=24,b=True,c=LBLUE)
    T(s,6,y+0.15,15,0.6,na,sz=13,b=True,c=NAVY)
    T(s,6,y+0.8,20,0.5,de,sz=9,c=GR)

# ═══ Slide 3: 研究背景 ═══
s=ns(); ST(s,1,'研究背景与意义','1/9'); FT(s)
# 左:背景
R(s,1,3,15.5,6.5,WH,LGR)
T(s,1.5,3.2,14.5,0.5,'研究背景',sz=12,b=True,c=NAVY)
T(s,1.5,3.8,14.5,5.5,'2024年网络零售额15.5万亿元，连续12年全球最大网络零售市场。电商规模持续扩大，仓储物流成为效率瓶颈。\n\nAMR凭借激光SLAM、视觉感知和自主决策能力，成为智能仓储核心装备。AGV行业规模从48.27亿升至95.49亿元，AMR市场突破60亿元，年增长率超30%。',sz=9.5,c=GR)
# 右:快递图(宽度约束,放大)
R(s,17,3,15.5,6.5,WH,LGR)
I(s,os.path.join(fig,'kuaidi1.png'),17.3,3.2,w=14.8,max_h=5.5)
T(s,17.3,8.9,14.8,0.4,'2020-2024年快递数据',sz=8,c=GR,a=PP_ALIGN.CENTER)
# 左下:AGV(宽度约束,放大)
R(s,1,10,10,6.2,WH,LGR)
I(s,os.path.join(fig,'AGV.png'),1.3,10.2,w=9.4,max_h=5.0)
T(s,1.3,15.3,9.4,0.4,'AGV工作示意图',sz=8,c=GR,a=PP_ALIGN.CENTER)
# 中下:目标
R(s,11.5,10,10.5,6.2,RGBColor(0xE8,0xF5,0xE9))
T(s,12,10.2,9.5,0.5,'研究目标',sz=11,b=True,c=RGBColor(0x2E,0x7D,0x32))
for i,g in enumerate(['① 自主沿规划路径行驶','② 载物台升降 ≥0.4m','③ 障碍物检测与避障','④ 紧凑外形，灵活通行','⑤ 结构简洁、成本可控']):
    T(s,12,10.8+i*0.55,9.5,0.5,g,sz=9,c=RGBColor(0x2E,0x7D,0x32))
# 右下:指标
R(s,22.5,10,10.3,6.2,WH,LGR)
T(s,23,10.2,9.3,0.5,'关键指标',sz=11,b=True,c=NAVY)
for i,(k,v) in enumerate([('额定载重','50 kg'),('升降行程','0.4 m'),('最大速度','30 m/min')]):
    y=10.8+i*1.6; T(s,23,y,4,0.4,k,sz=9,c=GR)
    T(s,23,y+0.35,9.3,0.7,v,sz=18,b=True,c=LBLUE)

# ═══ Slide 4: 机械结构 ═══
s=ns(); ST(s,2,'机械结构设计','2/9'); FT(s)
# 左:总装图(宽度约束)
R(s,1,3,10,7.8,WH,LGR)
I(s,os.path.join(fig,'机械总装.png'),1.3,3.2,w=9.4,max_h=6.5)
T(s,1.3,9.8,9.4,0.4,'总体结构示意图',sz=8,c=GR,a=PP_ALIGN.CENTER)
# 中:升降机构(宽度约束)
R(s,11.5,3,10,4.5,WH,LGR)
I(s,os.path.join(fig,'shengjiangtai.png'),11.8,3.2,w=9.4,max_h=3.5)
T(s,11.8,6.8,9.4,0.4,'剪叉式升降机构',sz=8,c=GR,a=PP_ALIGN.CENTER)
# 右:参数
R(s,22,3,10.5,4.5,RGBColor(0xE3,0xF2,0xFD))
T(s,22.5,3.2,9.5,0.5,'升降机构参数',sz=10,b=True,c=RGBColor(0x15,0x65,0xC0))
for i,p in enumerate(['臂杆 L=600mm','截面 36×4.8mm Q235','夹角 19.1°~133.5°','高度 149~600mm','升程 ~452mm','载重 50kg']):
    T(s,22.5,3.8+i*0.38,9.5,0.35,p,sz=8.5,c=RGBColor(0x15,0x65,0xC0))
# 中:FEA文字
R(s,11.5,8,21,4,RGBColor(0xE8,0xF5,0xE9))
T(s,12,8.2,20,0.5,'强度校核 (FEA验证)',sz=10,b=True,c=RGBColor(0x2E,0x7D,0x32))
for i,f in enumerate(['臂杆组合应力 5.56 MPa (FEA: 5.60, 偏差<1%)','Q235许用应力 156.7 MPa  利用率 3.5%','平台 von Mises 14.1 MPa  利用率 9.0%','平台挠度 0.47mm  |  屈曲安全系数 10.14']):
    T(s,12,8.8+i*0.48,20,0.42,f,sz=9,c=RGBColor(0x2E,0x7D,0x32))
# 底:FEA图(宽度约束,放大)
R(s,1,12.5,10.5,5.2,WH,LGR)
I(s,os.path.join(fig,'fea_results_50kg.png'),1.3,12.7,w=10,max_h=4.0)
T(s,1.3,16.8,10,0.4,'有限元分析结果',sz=7,c=GR,a=PP_ALIGN.CENTER)
# 底右:结论
R(s,12,12.5,20.5,5.2,RGBColor(0xFF,0xF8,0xE1))
T(s,12.5,12.7,19.5,0.5,'结论',sz=10,b=True,c=RGBColor(0xF5,0x7F,0x17))
T(s,12.5,13.3,19.5,4,'各部件应力远低于许用值，结构安全可靠。\n理论手算与有限元偏差<1%，相互印证。\n承载平台应力利用率仅9.0%，安全裕度充足。\n增设侧向支撑后屈曲安全系数提升至10.14。',sz=9.5,c=RGBColor(0xF5,0x7F,0x17))

# ═══ Slide 5: 电气系统 ═══
s=ns(); ST(s,3,'电气系统设计','3/9'); FT(s)
cols=[(1,'STM32F103C8T6',os.path.join(fig,'STM32F103C8T6实物图.png'),'ARM Cortex-M3\n72MHz / 64KB Flash\n37 I/O'),
      (9,'SMC80S 伺服电机',os.path.join(fig,'SMC80S-0040-30AoK-3DKH电机外形图.png'),'400W / 1.27N·m\n3000rpm / 48VDC\n2500P/R编码器'),
      (17,'TCRT5000 红外',os.path.join(fig,'TCRT5000.jpg'),'反射式光电传感器\n检测距离 1~25mm\n数字量输出'),
      (25,'HC-SR04 超声波',os.path.join(fig,'HC-SR04.jpg'),'声波飞行时间测距\n2cm~400cm\n精度3mm')]
for l,nm,pa,de in cols:
    R(s,l,3,7.5,7.5,WH,LGR)
    I(s,pa,l+0.3,3.2,w=6.9,max_h=5.0)
    T(s,l+0.3,8.5,6.9,0.4,nm,sz=9,b=True,c=NAVY,a=PP_ALIGN.CENTER)
    T(s,l+0.3,9,6.9,1.2,de,sz=8,c=GR,a=PP_ALIGN.CENTER)
# 下左:参数
R(s,1,11,15.5,5.5,RGBColor(0xE3,0xF2,0xFD))
T(s,1.5,11.2,14.5,0.5,'电气系统参数',sz=10,b=True,c=RGBColor(0x15,0x65,0xC0))
for i,e in enumerate(['驱动轮直径 180mm，轮距 520mm','减速器 ZJPX115 i=40，输出扭矩 47.2 N·m','驱动所需扭矩 22.5 N·m，功率余量充足','编码器 10000脉冲/转，轮速检测精度高']):
    T(s,1.5,11.8+i*0.55,14.5,0.5,e,sz=9,c=RGBColor(0x15,0x65,0xC0))
# 下右:传感器
R(s,17,11,15.5,5.5,WH,LGR)
I(s,os.path.join(fig,'sensor_layout.png'),17.3,11.2,w=6.5)
T(s,24.5,11.5,7.5,5,'5路红外传感器布局：\n间距40mm，覆盖160mm\n安装在车体前端80mm处\n权重{-4,-2,0,+2,+4}\n\n控制架构：\nSTM32采集→PID运算→PWM输出',sz=9,c=GR)

# ═══ Slide 6: 运动学与轨迹 ═══
s=ns(); ST(s,4,'运动学建模与轨迹仿真','4/9'); FT(s)
# 左:模型图(宽度约束)
R(s,1,3,14,7.5,WH,LGR)
I(s,os.path.join(fig,'双轮模型.png'),1.3,3.2,w=13.4,max_h=6.5)
T(s,1.3,9.9,13.4,0.4,'双轮差速模型',sz=8,c=GR,a=PP_ALIGN.CENTER)
# 左下:方程
R(s,1,11,14,5.5,RGBColor(0xE3,0xF2,0xFD))
T(s,1.5,11.2,13,0.5,'运动学方程',sz=10,b=True,c=RGBColor(0x15,0x65,0xC0))
T(s,1.5,11.8,13,4.5,'v = (v_R + v_L) / 2\nω = (v_R - v_L) / L\n\nẋ = v·cosθ    ẏ = v·sinθ    θ̇ = ω\n\n轮距 L=0.52m  轮半径 r=0.09m',sz=10,c=RGBColor(0x15,0x65,0xC0))
# 右:4轨迹(宽度约束)
traj=[(os.path.join(fig,'直线.png'),'直线',15.5,3),(os.path.join(fig,'圆形.png'),'圆弧',24,3),(os.path.join(fig,'s型.png'),'S形',15.5,9.5),(os.path.join(fig,'8字.png'),'8字',24,9.5)]
for pa,nm,l,t in traj:
    R(s,l,t,8,5.8,WH,LGR)
    I(s,pa,l+0.2,t+0.2,w=7.6,max_h=4.8)
    T(s,l+0.2,t+5,7.6,0.4,nm,sz=8.5,b=True,c=NAVY,a=PP_ALIGN.CENTER)
# 右下:结果表
R(s,15.5,15.8,17,2,WH,LGR)
T(s,16,15.9,16,0.4,'轨迹仿真结果汇总',sz=10,b=True,c=NAVY)
T(s,16,16.4,16,1.2,'直线: 路径9.0m, 位移9.0m, ω=0    |    圆弧: 路径7.0m, ω=0.45rad/s\nS形: 路径8.0m, ω=0.60rad/s, κ=1.5  |    8字: 路径7.6m, ω=0.70rad/s, κ=1.84',sz=8,c=GR)

# ═══ Slide 7: 路径规划 ═══
s=ns(); ST(s,5,'路径规划与循迹仿真','5/9'); FT(s)
# 上左:地图(宽度约束,放大)
R(s,1,3,20,7.5,WH,LGR)
I(s,os.path.join(mlb,'hybrid_astar_result.png'),1.3,3.2,w=12,max_h=6.5)
T(s,13.5,3.4,7,0.5,'Hybrid A* 路径规划',sz=10,b=True,c=NAVY)
T(s,13.5,4,7,5.5,'算法：节点(x,y,θ)三维位姿\n5种曲率运动基元扩展\n代价 f=g+h\n\n参数：\n地图 10m×8m, 0.2m\n5个货架\n起点(1,1)→终点(9,7)\n迭代3008次\n路径11.04m, 25点',sz=8.5,c=GR)
# 上右:指标
R(s,21.5,3,11,7.5,RGBColor(0xE8,0xF5,0xE9))
T(s,22,3.2,10,0.5,'规划结果',sz=10,b=True,c=RGBColor(0x2E,0x7D,0x32))
for i,(k,v) in enumerate([('路径长度','11.04 m'),('路径点数','25'),('搜索迭代','3008'),('曲率κ','±1.11, ±0.52, 0'),('步长','0.5 m'),('轮距L','0.52 m')]):
    T(s,22,3.8+i*0.6,5,0.5,k,sz=9,c=GR);T(s,27,3.8+i*0.6,5,0.5,v,sz=9,b=True,c=RGBColor(0x2E,0x7D,0x32))
# 下左:轨迹(宽度约束)
R(s,1,11,15.5,5.5,WH,LGR)
T(s,1.5,11.2,14.5,0.5,'参考路径 vs 循迹轨迹',sz=10,b=True,c=NAVY)
I(s,os.path.join(mlb,'line_following_traj.png'),1.3,11.8,w=14.8,max_h=4.5)
# 下右:误差(宽度约束)
R(s,17,11,15.5,5.5,WH,LGR)
T(s,17.5,11.2,14.5,0.5,'循迹横向误差',sz=10,b=True,c=NAVY)
I(s,os.path.join(mlb,'line_following_error.png'),17.3,11.8,w=14.8,max_h=4.5)
# 底部指标
for i,(k,v,bg,tc) in enumerate([('RMS误差','1.5cm',RGBColor(0xE8,0xF5,0xE9),RGBColor(0x2E,0x7D,0x32)),('最大误差','6.8cm',RGBColor(0xFF,0xF3,0xE0),RGBColor(0xE6,0x51,0x00)),('路径长度','11.04m',RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0)),('耗时','36.2s',RGBColor(0xF3,0xE5,0xF5),RGBColor(0x7B,0x1F,0xA2))]):
    cx=1+i*8.2; R(s,cx,17.2,7.5,0.9,bg,tc)
    T(s,cx+0.3,17.25,3,0.35,k,sz=8,c=tc)
    T(s,cx+3.5,17.25,3.5,0.35,v,sz=10,b=True,c=tc,a=PP_ALIGN.RIGHT)

# ═══ Slide 8: 运动基元 ═══
s=ns(); ST(s,6,'运动基元与算法流程','6/9'); FT(s)
# 左:5基元
R(s,1,3,15,7.5,WH,LGR)
T(s,1.5,3.2,14,0.5,'差速运动基元 (5种曲率)',sz=10,b=True,c=NAVY)
cls=[RGBColor(0xE6,0x51,0x00),RGBColor(0xF3,0x9C,0x12),RGBColor(0x27,0xAE,0x60),RGBColor(0x29,0x80,0xB9),RGBColor(0x8E,0x44,0xAD)]
nms=['大左转','小左转','直行','小右转','大右转']
ks=['κ=-1.11','κ=-0.52','κ=0','κ=+0.52','κ=+1.11']
rs=['R=0.90m','R=1.92m','R=∞','R=1.92m','R=0.90m']
for i in range(5):
    y=3.9+i*0.65; R(s,1.5,y,14,0.58,RGBColor(0xFA,0xFA,0xFA),cls[i])
    T(s,2,y+0.08,3.5,0.38,nms[i],sz=9,b=True,c=cls[i])
    T(s,5.5,y+0.08,4,0.38,ks[i],sz=8.5,c=GR)
    T(s,9.5,y+0.08,4,0.38,rs[i],sz=8.5,c=GR)
T(s,1.5,7.5,14,2.5,'前进/后退 × 5种曲率 = 10种基元\n步长 d=0.5m\n节点状态 (x,y,θ) 三维位姿',sz=9,c=GR)
# 右:流程
R(s,16.5,3,16,7.5,WH,LGR)
T(s,17,3.2,15,0.5,'Hybrid A* 算法流程',sz=10,b=True,c=NAVY)
steps=[('1. 栅格地图','10×8m, 0.2m'),('2. 初始化','Open/Close, h₀'),('3. 节点扩展','10种基元'),('4. 碰撞检测','采样11点'),('5. 代价评估','f=g+h'),('6. 路径回溯','父节点链'),('7. 均值平滑','3次滤波'),('8. 等距插值','0.05m, 224点')]
for i,(nm,de) in enumerate(steps):
    y=3.8+i*0.52; bg=RGBColor(0xE8,0xEA,0xF6) if i<5 else RGBColor(0xE8,0xF5,0xE9)
    R(s,17,y,15,0.46,bg)
    T(s,17.5,y+0.04,5.5,0.38,nm,sz=8,b=True,c=NAVY)
    T(s,23,y+0.04,8,0.38,de,sz=8,c=GR)
# 下:代价函数
R(s,1,11,31.5,2.5,RGBColor(0xFF,0xF8,0xE1))
T(s,1.5,11.2,30,0.5,'代价函数 f = g + h',sz=10,b=True,c=RGBColor(0xF5,0x7F,0x17))
T(s,1.5,11.8,10,1.5,'g: 实际代价 (步长0.5m\n+ 倒车+50% + 转弯+0.1)',sz=9,c=RGBColor(0xF5,0x7F,0x17))
T(s,12,11.8,10,1.5,'h: 欧氏距离启发值\n(保证可采纳性)',sz=9,c=RGBColor(0xF5,0x7F,0x17))
T(s,23,11.8,9,1.5,'关键：节点含θ\n满足差速约束\n路径可直接执行',sz=9,b=True,c=RGBColor(0xF5,0x7F,0x17))
# 底:传感器
R(s,1,14,15.5,3.8,WH,LGR)
I(s,os.path.join(fig,'sensor_layout.png'),1.3,14.2,w=7,max_h=3.5)
T(s,9,14.5,7,3,'循迹传感器：\n5路TCRT5000, 间距40mm\n权重{-4,-2,0,+2,+4}\nPID: Kp=3.0 Ki=0.08 Kd=1.0',sz=9,c=GR)
R(s,17,14,15.5,3.8,RGBColor(0xE8,0xF5,0xE9))
T(s,17.5,14.2,14.5,0.5,'两层导航架构',sz=10,b=True,c=RGBColor(0x2E,0x7D,0x32))
T(s,17.5,14.8,14.5,2.5,'规划层 Hybrid A* → 确定"走哪条路"\n执行层 PID+红外 → 解决"如何沿路走"\n\n优势：地图变更后自动重规划',sz=9,c=RGBColor(0x2E,0x7D,0x32))

# ═══ Slide 9: 双环PID ═══
s=ns(); ST(s,7,'双环PID运动控制','7/9'); FT(s)
# 左:结构图(宽度约束)
R(s,1,3,12,8,WH,LGR)
I(s,os.path.join(mlb,'cascaded_pid_diagram.png'),1.3,3.2,w=11.4,max_h=5.0)
T(s,1.3,8.5,11.4,0.4,'双环PID级联控制结构',sz=8,c=GR,a=PP_ALIGN.CENTER)
# 左下:参数
R(s,1,11.5,12,5.5,RGBColor(0xE3,0xF2,0xFD))
T(s,1.5,11.8,11,0.5,'PID参数',sz=10,b=True,c=RGBColor(0x15,0x65,0xC0))
T(s,1.5,12.4,11,0.4,'位姿环(航向PID):  Kp=3.0  Ki=0.2  Kd=0.4',sz=9,c=RGBColor(0x15,0x65,0xC0))
T(s,1.5,12.9,11,0.4,'轮速环(PI):  Kp=17.0  Ki=103.0',sz=9,c=RGBColor(0x15,0x65,0xC0))
T(s,1.5,13.6,11,3,'控制策略：\n• 内环响应频率约为外环3~5倍\n• 轮速环用PI不用PID(微分对噪声敏感)\n• 控制周期 5ms，中断方式执行\n• 积分限幅[-2,2]防饱和\n• 电压限幅 ±24V',sz=9,c=GR)
# 右:4仿真图(宽度约束)
pid=[(os.path.join(mlb,'pid_step_response.png'),'阶跃响应',13.5,3),(os.path.join(mlb,'pid_point_stabilization.png'),'定点镇定',23,3),(os.path.join(mlb,'pid_line_tracking.png'),'直线跟踪',13.5,9.5),(os.path.join(mlb,'pid_circle_tracking.png'),'圆形跟踪',23,9.5)]
for pa,nm,l,t in pid:
    R(s,l,t,9,5.8,WH,LGR)
    I(s,pa,l+0.2,t+0.2,w=8.6,max_h=5.0)
    T(s,l+0.2,t+5.3,8.6,0.4,nm,sz=8.5,b=True,c=NAVY,a=PP_ALIGN.CENTER)
# 底:结果
R(s,1,15.8,31.5,2,WH,LGR)
T(s,1.5,15.9,30,0.5,'仿真结果汇总',sz=10,b=True,c=NAVY)
T(s,1.5,16.5,30,1,'阶跃: 稳态误差≈0, 调节0.8s  |  定点: 误差(0.022,-0.007)m, 航向-0.89°  |  直线: 角速度RMSE 0.0023rad/s  |  圆形: 线速度RMSE 0.0086m/s, 圆度RMSE 0.020m',sz=8.5,c=GR)

# ═══ Slide 10: 致谢 ═══
s=ns(); R(s,0,0,SW,SH,NAVY)
R(s,2,6,8,0.04,LBLUE)
T(s,2,6.5,30,1.5,'感谢各位老师指导',sz=30,b=True,c=WH)
R(s,2,8.5,8,0.04,LBLUE)
T(s,2,9.5,30,0.5,'答辩人：王梓煜    指导教师：曹树春',sz=13,c=RGBColor(0xBB,0xDE,0xFB))
T(s,2,10.3,30,0.5,'齐鲁工业大学  机械工程学院',sz=13,c=RGBColor(0xBB,0xDE,0xFB))

# 保存
prs.save(new_pptx)
print(f"已保存: {new_pptx}")
print(f"共 {len(prs.slides)} 页, 尺寸 {prs.slide_width/360000:.1f}x{prs.slide_height/360000:.1f}cm")
