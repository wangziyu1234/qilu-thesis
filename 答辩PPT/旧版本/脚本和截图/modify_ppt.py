# -*- coding: utf-8 -*-
"""
为答辩PPT添加论文引用的图片，实现图文并茂
只使用论文中通过 \includegraphics 引用的图片
策略：在现有文字内容之间的空隙区域添加图片，不覆盖原有文字
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

IMG_BASE = r"D:\bylw\code\lunwen\QLULatex\QLUThesisLatexTemplate-master\Thesis\static\figures"

def img(name):
    return os.path.join(IMG_BASE, name)

def add_img(slide, img_path, left, top, width=None, height=None):
    if not os.path.exists(img_path):
        print(f"  [WARN] 不存在: {img_path}")
        return None
    try:
        if width and height:
            pic = slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            pic = slide.shapes.add_picture(img_path, left, top)
        print(f"  [OK] {os.path.basename(img_path)}")
        return pic
    except Exception as e:
        print(f"  [ERR] {os.path.basename(img_path)}: {e}")
        return None

def add_label(slide, text, left, top, width, size=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Pt(14))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p.alignment = PP_ALIGN.CENTER
    return txBox

def main():
    prs = Presentation(r"D:\bylw\code\答辩PPT\答辩PPT_完整版.pptx")
    slides = list(prs.slides)
    print(f"总幻灯片数: {len(slides)}\n")

    # ================================================================
    # Slide 3: 01 研究背景与意义
    # 三列文字(2.9-10.7) + 下方通栏(12.9-18.1)均填满，无空隙
    # ================================================================
    print("=== Slide 3: 01 研究背景与意义 ===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 4: 02 研究目标与任务
    # 指标卡片(3.0-5.4)和下方文字(9.3-18.1)均填满，无空隙
    # ================================================================
    print("\n=== Slide 4: 02 研究目标与任务 ===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 5: 03 总体设计方案（驱动与升降）
    # Text 14 "选型依据"(1.4-32.4, 13.5-14.2)
    # Text 16 (1.1-32.1, 15.9-16.4) 单行文字
    # 14.2-15.9 之间有约1.7cm空隙，可放小图
    # ================================================================
    print("\n=== Slide 5: 03 总体设计方案（驱动与升降）===")
    s = slides[4]
    # 14.2-15.9空隙中添加驱动方式小图
    add_img(s, img("chasu.png"), Cm(2.0), Cm(14.3), width=Cm(3.5))
    add_img(s, img("danlunqudong.png"), Cm(6.0), Cm(14.3), width=Cm(3.5))
    add_img(s, img("quanfangwei.png"), Cm(10.0), Cm(14.3), width=Cm(3.5))
    add_label(s, "左:差速驱动  中:单轮驱动  右:全方位驱动（论文图2-1~2-3）", Cm(1.0), Cm(16.5), Cm(14.0))

    # ================================================================
    # Slide 6: 04 总体设计方案（控制与规划）
    # 文字区: 左上(1.4-15.9,4.0-9.0) 右上(17.1-32.4,4.0-9.0)
    #         左下(1.4-15.9,11.1-18.1) 右下(17.1-32.4,11.1-18.1)
    # 策略: 文字填满了整个区域，空间极其有限，跳过
    # ================================================================
    print("\n=== Slide 6: 04 总体设计方案（控制与规划）===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 7: 05 机械结构设计（剪叉机构）
    # 文字区: 左上(1.4-14.6,4.0-10.2) 右上(15.9-32.4,4.0-10.2)
    #         左下(1.4-15.3,12.9-15.9) 右下(15.9-32.4,13.6-16.7)
    # 左下文字到15.9，形状到18.3，有2.4cm空隙
    # ================================================================
    print("\n=== Slide 7: 05 机械结构设计（剪叉机构）===")
    s = slides[6]
    # 左下材料选择文字下方(15.9-18.3)添加运动简图
    add_img(s, img("jiagou_low.png"), Cm(2.0), Cm(16.1), height=Cm(2.0))
    add_img(s, img("jiagou_high.png"), Cm(6.5), Cm(16.1), height=Cm(2.0))
    add_label(s, "左:最低位  右:最高位（论文图3-2/3-3）", Cm(2.0), Cm(18.2), Cm(12.0))

    # ================================================================
    # Slide 8: 06 机械结构设计（强度校核）
    # 文字区到16.5，形状到18.3，仅1.8cm空隙，空间不足
    # ================================================================
    print("\n=== Slide 8: 06 机械结构设计（强度校核）===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 9: 07 电气系统设计（控制器与电机）
    # 上下文字区之间仅1.1cm空隙，空间不足
    # ================================================================
    print("\n=== Slide 9: 07 电气系统设计（控制器与电机）===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 10: 08 电气系统设计（传感器与电路）
    # 上下文字区之间仅1.0cm空隙，空间不足
    # ================================================================
    print("\n=== Slide 10: 08 电气系统设计（传感器与电路）===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 11: 09 运动学建模与仿真
    # 左列: 模型(1.4-14.6,4.0-9.7) + 参数(1.4-14.6,11.9-16.8)
    # 右列: 标题(15.5-32.0,2.8-3.6) + 表格(15.5-32.8,3.6-7.2) + 结论(15.9-32.4,11.9-16.8)
    # 策略: 在左列"仿真参数"下方和右列"仿真结论"下方添加图
    # ================================================================
    print("\n=== Slide 11: 09 运动学建模与仿真 ===")
    s = slides[10]
    # 左列模型区域(4.0-9.7)底部添加双轮模型图
    add_img(s, img("双轮模型.png"), Cm(2.0), Cm(7.5), width=Cm(6.5))
    add_label(s, "双轮差速机器人模型（论文图5-1）", Cm(2.0), Cm(9.8), Cm(6.5))
    # 右列"四类典型轨迹"表格下方(7.2-10.7有约3.5cm空隙)
    add_img(s, img("直线.png"), Cm(16.0), Cm(7.5), width=Cm(3.8))
    add_label(s, "直线", Cm(16.0), Cm(10.2), Cm(3.8))
    add_img(s, img("圆形.png"), Cm(20.0), Cm(7.5), width=Cm(3.8))
    add_label(s, "圆弧", Cm(20.0), Cm(10.2), Cm(3.8))
    add_img(s, img("s型.png"), Cm(24.0), Cm(7.5), width=Cm(3.8))
    add_label(s, "S形", Cm(24.0), Cm(10.2), Cm(3.8))
    add_img(s, img("8字.png"), Cm(28.0), Cm(7.5), width=Cm(3.8))
    add_label(s, "8字", Cm(28.0), Cm(10.2), Cm(3.8))

    # ================================================================
    # Slide 12: 10 路径规划与循迹
    # 上: 左(1.4-15.9,4.0-9.7) 右(17.1-32.4,4.0-9.7)
    # 中: 通栏标题(1.4-32.4,10.7-11.6) + 文字(1.4-32.4,11.9-14.3)
    # 下: 通栏标题(1.4-32.4,15.3-16.2) + 文字(1.3-32.4,15.5-18.5)
    # 上下文字区之间(9.7-10.7)有1cm空隙
    # 中间文字区和下方文字区之间(14.3-15.3)有1cm空隙
    # ================================================================
    print("\n=== Slide 12: 10 路径规划与循迹 ===")
    s = slides[11]
    # 上方两列区域(4.0-9.7)内，文字不填满全部空间
    # 在右列文字区底部(约8.5-9.7)添加传感器布局图
    add_img(s, img("sensor_layout.png"), Cm(17.5), Cm(8.0), width=Cm(4.0))
    add_label(s, "传感器布局（论文图5-8）", Cm(17.5), Cm(10.0), Cm(4.0))
    # 上下区域之间的1cm空隙(9.7-10.7)放小图
    add_img(s, img("hybrid_astar_result.png"), Cm(2.0), Cm(9.8), width=Cm(6.5))
    add_img(s, img("line_following_traj.png"), Cm(9.0), Cm(9.8), width=Cm(6.5))
    add_label(s, "左:Hybrid A*（论文图5-7）  右:循迹轨迹（论文图5-9）", Cm(2.0), Cm(10.8), Cm(14.0))
    # 中间文字区和下方文字区之间的1cm空隙(14.3-15.3)放误差图
    add_img(s, img("line_following_error.png"), Cm(2.0), Cm(14.4), width=Cm(6.5))
    add_img(s, img("line_following_omega.png"), Cm(9.0), Cm(14.4), width=Cm(6.5))
    add_label(s, "左:横向误差（论文图5-10）  右:纠偏角速度（论文图5-11）", Cm(2.0), Cm(15.3), Cm(14.0))

    # ================================================================
    # Slide 13: 11 运动控制算法设计
    # 上下文字区之间仅1.1cm空隙，空间不足
    # ================================================================
    print("\n=== Slide 13: 11 运动控制算法设计 ===")
    print("  跳过（文字区域已满）")

    # ================================================================
    # Slide 14: 12 运动控制仿真验证
    # 左上: 表格(1.0-16.3, 3.6-8.6) + 标题(1.0-15.0, 2.8-3.6)
    # 右上: 阶跃响应文字(17.1-32.4, 4.0-10.2)
    # 左下: 直线与圆形跟踪文字(1.4-15.9, 12.4-18.1)
    # 右下: 定点镇定文字(17.1-32.4, 12.4-18.1)
    # 策略: 在表格下方(8.6-11.2有2.6cm空隙)添加图
    # ================================================================
    print("\n=== Slide 14: 12 运动控制仿真验证 ===")
    s = slides[13]
    # 表格下方空隙(8.6-11.2)添加四场景图
    add_img(s, img("pid_step_response.png"), Cm(1.0), Cm(8.8), width=Cm(3.5))
    add_label(s, "阶跃", Cm(1.0), Cm(11.0), Cm(3.5))
    add_img(s, img("pid_point_stabilization.png"), Cm(4.8), Cm(8.8), width=Cm(3.5))
    add_label(s, "镇定", Cm(4.8), Cm(11.0), Cm(3.5))
    add_img(s, img("pid_line_tracking.png"), Cm(8.6), Cm(8.8), width=Cm(3.5))
    add_label(s, "直线", Cm(8.6), Cm(11.0), Cm(3.5))
    add_img(s, img("pid_circle_tracking.png"), Cm(12.4), Cm(8.8), width=Cm(3.5))
    add_label(s, "圆形", Cm(12.4), Cm(11.0), Cm(3.5))

    # ================================================================
    # Slide 15: 13 结论与展望
    # 左列(1.4-15.9,4.0-15.3) 右列(17.1-32.4,4.0-15.3)
    # 右列文字到15.3，下方形状到15.5，再下方"应用前景"从16.3开始
    # 15.5-16.3之间仅0.8cm空隙，空间不足
    # ================================================================
    print("\n=== Slide 15: 13 结论与展望 ===")
    print("  跳过（文字区域已满）")

    # 保存
    output = r"D:\bylw\code\答辩PPT\答辩PPT_图文版.pptx"
    prs.save(output)
    print(f"\n{'='*50}")
    print(f"完成！已保存到: {output}")
    print(f"{'='*50}")

if __name__ == "__main__":
    main()
