"""答辩PPT完整重排：11页，新建文件"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os, shutil

old_pptx = 'D:/bylw/code/答辩PPT/答辩PPT_图文并茂.pptx'
new_pptx = 'D:/bylw/code/答辩PPT/答辩PPT_重排版.pptx'
fig = 'D:/bylw/code/lunwen/QLULatex/QLUThesisLatexTemplate-master/Thesis/static/figures'
mlb = 'D:/bylw/code/fangzhen/matlab'

# 读取旧PPT获取封面和致谢
old_prs = Presentation(old_pptx)

# 创建新PPT (使用旧PPT的模板)
prs = Presentation(old_pptx)
# 删除所有幻灯片
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

SW = prs.slide_width; SH = prs.slide_height

DB = RGBColor(0x28,0x35,0x93); WH = RGBColor(0xFF,0xFF,0xFF)
GR = RGBColor(0x66,0x66,0x66); LT = RGBColor(0x55,0x55,0x55)

def rrect(sl,l,t,w,h,fill,line=None):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Cm(l),Cm(t),Cm(w),Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line: s.line.color.rgb=line; s.line.width=Pt(1)
    else: s.line.fill.background()
    return s

def txt(sl,l,t,w,h,text,sz=12,bold=False,color=GR,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Cm(l),Cm(t),Cm(w),Cm(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text
    p.font.size=Pt(sz); p.font.bold=bold; p.font.color.rgb=color; p.alignment=align

def img(sl,path,l,t,w=None,h=None):
    if w and h: return sl.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w),height=Cm(h))
    elif w: return sl.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w))
    elif h: return sl.shapes.add_picture(path,Cm(l),Cm(t),height=Cm(h))
    else: return sl.shapes.add_picture(path,Cm(l),Cm(t))

def title_bar(sl, num, name, page):
    rrect(sl,0,0,33.9,1.4,DB)
    rrect(sl,0,1.4,33.9,0.2,RGBColor(0x1A,0x23,0x7E))
    rrect(sl,0.7,0.1,2.0,1.1,RGBColor(0x3F,0x51,0xB5))
    txt(sl,0.6,0.4,2.0,0.7,num,sz=22,bold=True,color=WH,align=PP_ALIGN.CENTER)
    txt(sl,3.0,0.4,23,0.7,name,sz=18,bold=True,color=WH)
    txt(sl,30.5,2.3,2.2,0.4,page,sz=10,color=RGBColor(0xBB,0xBB,0xBB),align=PP_ALIGN.RIGHT)

def footer(sl):
    txt(sl,1.0,18.5,20,0.4,'齐鲁工业大学  机械工程学院',sz=8,color=RGBColor(0x99,0x99,0x99))

def new_slide():
    for i, layout in enumerate(prs.slide_layouts):
        if layout.name == 'Blank' or layout.placeholders == 0:
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[0])

# ══════════════════════════════════════════
# 从旧PPT复制封面
# ══════════════════════════════════════════
# 由于python-pptx不支持直接复制幻灯片，我们手动重建封面
s1 = new_slide()
rrect(s1,0,0,33.9,19.1,RGBColor(0x1A,0x23,0x7E))
rrect(s1,0,0,33.9,3.5,RGBColor(0x28,0x35,0x93))
txt(s1,2,4.5,30,1.5,'智能仓储移动机器人',sz=32,bold=True,color=WH,align=PP_ALIGN.CENTER)
txt(s1,2,6.0,30,1.0,'结构设计与运动控制研究',sz=24,color=RGBColor(0xBB,0xDE,0xFB),align=PP_ALIGN.CENTER)
txt(s1,2,7.5,30,0.8,'Design and Control of Intelligent Warehouse Mobile Robot',sz=14,color=RGBColor(0x90,0xCA,0xF9),align=PP_ALIGN.CENTER)
rrect(s1,8,9.5,18,0.03,RGBColor(0x3F,0x51,0xB5))
txt(s1,2,10.5,30,0.6,'答辩人：王梓煜    学号：202201230042    指导教师：曹树春',sz=12,color=RGBColor(0xBB,0xDE,0xFB),align=PP_ALIGN.CENTER)
txt(s1,2,11.3,30,0.6,'齐鲁工业大学  机械工程学院',sz=12,color=RGBColor(0xBB,0xDE,0xFB),align=PP_ALIGN.CENTER)
txt(s1,2,12.5,30,0.6,'2026年6月',sz=14,color=WH,align=PP_ALIGN.CENTER)
# 尝试添加logo
logo_path = os.path.join(fig,'MainLOGO-pretty.png')
if os.path.exists(logo_path):
    try: img(s1, logo_path, 14, 13.5, h=3.0)
    except: pass
print("Slide 1: 封面")

# ══════════════════════════════════════════
# 目录页 (简化重建)
# ══════════════════════════════════════════
s2 = new_slide()
rrect(s2,0,0,33.9,19.1,WH)
rrect(s2,0,0,33.9,1.6,DB)
txt(s2,1.0,0.3,5,1.0,'目  录',sz=28,bold=True,color=WH)
txt(s2,30,0.3,3,1.0,'CONTENTS',sz=12,color=RGBColor(0xBB,0xDE,0xFB),align=PP_ALIGN.RIGHT)

toc_items = [
    ('01','研究背景与意义','研究现状、目标与方法'),
    ('02','机械结构设计','总装布局、升降机构、强度校核'),
    ('03','电气系统设计','主控、电机、传感器、电路'),
    ('04','运动学建模','差速模型、速度与位姿关系'),
    ('05','轨迹仿真验证','直线/圆弧/S形/8字轨迹'),
    ('06','路径规划与循迹','Hybrid A* + PID红外循迹'),
    ('07','运动基元与算法','5种曲率基元、算法流程'),
    ('08','双环PID控制','级联结构、四项仿真验证'),
]
for i,(num,name,desc) in enumerate(toc_items):
    y = 2.5 + i * 2.0
    rrect(s2, 2, y, 30, 1.7, RGBColor(0xF5,0xF5,0xF5))
    rrect(s2, 2, y, 0.15, 1.7, DB)
    txt(s2, 3, y+0.2, 2, 1.0, num, sz=28, bold=True, color=DB)
    txt(s2, 6, y+0.2, 12, 0.7, name, sz=14, bold=True, color=RGBColor(0x33,0x33,0x33))
    txt(s2, 6, y+0.9, 20, 0.6, desc, sz=10, color=GR)
print("Slide 2: 目录")

# ══════════════════════════════════════════
# Slide 3: 研究背景
# ══════════════════════════════════════════
s3 = new_slide()
title_bar(s3,'01','研究背景与意义','1/11')
rrect(s3,1.0,2.0,15.5,7.5,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s3,1.5,2.2,14.5,0.6,'研究背景',sz=12,bold=True,color=DB)
txt(s3,1.5,2.9,14.5,6.0,
    '2024年网络零售额15.5万亿元，连续12年全球最大网络零售市场。'
    '电商规模持续扩大，仓储物流环节成为效率瓶颈。'
    '自主移动机器人(AMR)凭借激光SLAM、视觉感知和自主决策能力，'
    '成为智能仓储的核心装备。'
    '\n\n'
    'AGV行业规模从48.27亿攀升至95.49亿元，'
    'AMR市场突破60亿元，年增长率超30%。'
    '行业从"固定路径导引"向"自主感知决策"升级。',
    sz=9,color=GR)
img(s3,os.path.join(fig,'kuaidi1.png'),17.0,2.0,w=15.5)
txt(s3,17.0,6.3,15.5,0.4,'图  2020-2024年人均用邮与快递情况',sz=7,color=LT,align=PP_ALIGN.CENTER)
img(s3,os.path.join(fig,'AGV.png'),1.0,10.0,w=8.0)
txt(s3,1.0,14.3,8.0,0.4,'AGV工作示意图',sz=7,color=LT,align=PP_ALIGN.CENTER)
rrect(s3,17.0,10.0,15.5,6.5,RGBColor(0xE8,0xF5,0xE9),RGBColor(0x43,0xA0,0x47))
txt(s3,17.5,10.2,14.5,0.6,'研究目标',sz=11,bold=True,color=RGBColor(0x2E,0x7D,0x32))
for i,g in enumerate(['① 自主沿规划路径行驶，完成物料搬运','② 载物台升降，行程≥0.4m','③ 基本障碍物检测与避障','④ 外形紧凑，狭窄通道灵活通行','⑤ 结构简洁、控制可靠、成本可控']):
    txt(s3,17.5,10.9+i*0.55,14.5,0.5,g,sz=9,color=RGBColor(0x2E,0x7D,0x32))
for i,(lab,val,bg,tc) in enumerate([('额定载重','50 kg',RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0)),('升降行程','0.4 m',RGBColor(0xE8,0xF5,0xE9),RGBColor(0x2E,0x7D,0x32)),('最大速度','30 m/min',RGBColor(0xFF,0xF3,0xE0),RGBColor(0xE6,0x51,0x00))]):
    cx=1.0+i*5.5
    rrect(s3,cx,15.0,5.0,1.5,bg,tc)
    txt(s3,cx+0.2,15.1,4.6,0.5,lab,sz=8,color=tc,align=PP_ALIGN.CENTER)
    txt(s3,cx+0.2,15.6,4.6,0.7,val,sz=18,bold=True,color=tc,align=PP_ALIGN.CENTER)
footer(s3)
print("Slide 3: 研究背景")

# ══════════════════════════════════════════
# Slide 4: 机械结构
# ══════════════════════════════════════════
s4 = new_slide()
title_bar(s4,'02','机械结构设计','2/11')
img(s4,os.path.join(fig,'机械总装.png'),1.0,2.0,w=11.0)
txt(s4,1.0,10.5,11.0,0.4,'智能仓储移动机器人总体结构',sz=8,color=LT,align=PP_ALIGN.CENTER)
img(s4,os.path.join(fig,'shengjiangtai.png'),12.5,2.0,w=7.0)
txt(s4,12.5,7.5,7.0,0.4,'剪叉式升降机构三维模型',sz=8,color=LT,align=PP_ALIGN.CENTER)
rrect(s4,20.0,2.0,12.5,5.5,RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0))
txt(s4,20.5,2.2,11.5,0.5,'升降机构参数',sz=10,bold=True,color=RGBColor(0x15,0x65,0xC0))
for i,p in enumerate(['臂杆长度 L = 600 mm','截面 36mm x 4.8mm Q235扁钢','夹角范围 19.1° ~ 133.5°','台面高度 149 ~ 600 mm','最大升程 ~452 mm','额定载重 50 kg']):
    txt(s4,20.5,2.8+i*0.45,11.5,0.4,p,sz=8,color=RGBColor(0x15,0x65,0xC0))
rrect(s4,20.0,8.0,12.5,4.5,RGBColor(0xE8,0xF5,0xE9),RGBColor(0x2E,0x7D,0x32))
txt(s4,20.5,8.2,11.5,0.5,'强度校核结果 (FEA验证)',sz=10,bold=True,color=RGBColor(0x2E,0x7D,0x32))
for i,f in enumerate(['臂杆组合应力  5.56 MPa  (FEA: 5.60)','许用应力  156.7 MPa  利用率 3.5%','平台 von Mises  14.1 MPa  利用率 9.0%','平台挠度  0.47 mm','屈曲安全系数  10.14 (含侧向支撑)']):
    txt(s4,20.5,8.8+i*0.45,11.5,0.4,f,sz=8,color=RGBColor(0x2E,0x7D,0x32))
img(s4,os.path.join(fig,'fea_results_50kg.png'),1.0,11.5,w=18.5)
txt(s4,1.0,16.5,18.5,0.4,'剪叉式升降机构有限元分析结果（50kg额定载重）',sz=7,color=LT,align=PP_ALIGN.CENTER)
txt(s4,20.0,13.0,12.5,3.0,'各部件应力远低于许用值，结构安全可靠。\n理论手算与有限元偏差<1%，相互印证。',sz=9,bold=True,color=RGBColor(0x2E,0x7D,0x32))
footer(s4)
print("Slide 4: 机械结构")

# ══════════════════════════════════════════
# Slide 5: 电气系统
# ══════════════════════════════════════════
s5 = new_slide()
title_bar(s5,'03','电气系统设计','3/11')
img(s5,os.path.join(fig,'STM32F103C8T6实物图.png'),1.0,2.0,w=5.5)
txt(s5,1.0,7.5,5.5,0.4,'STM32F103C8T6',sz=8,color=LT,align=PP_ALIGN.CENTER)
txt(s5,1.0,7.9,5.5,1.5,'ARM Cortex-M3\n72MHz / 64KB Flash\n20KB RAM / 37 I/O',sz=8,color=GR,align=PP_ALIGN.CENTER)
img(s5,os.path.join(fig,'SMC80S-0040-30AoK-3DKH电机外形图.png'),7.0,2.0,w=5.5)
txt(s5,7.0,7.5,5.5,0.4,'SMC80S 伺服电机',sz=8,color=LT,align=PP_ALIGN.CENTER)
txt(s5,7.0,7.9,5.5,1.5,'400W / 1.27N·m\n3000rpm / 48VDC\n2500P/R编码器',sz=8,color=GR,align=PP_ALIGN.CENTER)
img(s5,os.path.join(fig,'TCRT5000.jpg'),13.0,2.0,w=5.0)
txt(s5,13.0,6.5,5.0,0.4,'TCRT5000 红外',sz=8,color=LT,align=PP_ALIGN.CENTER)
img(s5,os.path.join(fig,'HC-SR04.jpg'),18.5,2.0,w=4.5)
txt(s5,18.5,6.5,4.5,0.4,'HC-SR04 超声波',sz=8,color=LT,align=PP_ALIGN.CENTER)
rrect(s5,24.0,2.0,8.5,6.0,RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0))
txt(s5,24.5,2.2,7.5,0.5,'电气系统参数',sz=10,bold=True,color=RGBColor(0x15,0x65,0xC0))
for i,e in enumerate(['驱动轮直径 180mm','轮距 520mm / 轮半径 90mm','减速器 ZJPX115 i=40','减速后扭矩 47.2 N·m','驱动所需扭矩 22.5 N·m','传感器: 5路红外+超声波']):
    txt(s5,24.5,2.8+i*0.45,7.5,0.4,e,sz=8,color=RGBColor(0x15,0x65,0xC0))
img(s5,os.path.join(fig,'最小系统原理图.png'),1.0,9.0,w=14.0)
txt(s5,1.0,14.8,14.0,0.4,'STM32F103C8T6最小系统原理图',sz=7,color=LT,align=PP_ALIGN.CENTER)
img(s5,os.path.join(fig,'sensor_layout.png'),15.5,9.0,w=6.5)
txt(s5,15.5,13.8,6.5,0.4,'5路红外传感器布局',sz=7,color=LT,align=PP_ALIGN.CENTER)
txt(s5,22.5,9.2,10.0,5.0,'控制架构：\nSTM32为中枢，采集传感器数据，\n执行PID算法，输出PWM驱动电机。\n\n传感器协同：\nTCRT5000近距离循迹(1~25mm)\nHC-SR04中远距离避障(2~400cm)\n编码器轮速反馈(10000脉冲/转)',sz=8,color=GR)
footer(s5)
print("Slide 5: 电气系统")

# ══════════════════════════════════════════
# Slide 6: 运动学建模
# ══════════════════════════════════════════
s6 = new_slide()
title_bar(s6,'04','运动学建模','4/11')
img(s6,os.path.join(fig,'双轮模型.png'),1.0,2.0,w=14.0)
txt(s6,1.0,10.5,14.0,0.4,'双轮差速移动机器人模型',sz=8,color=LT,align=PP_ALIGN.CENTER)
rrect(s6,15.5,2.0,17.0,5.0,RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0))
txt(s6,16.0,2.2,16.0,0.5,'差速运动学方程',sz=11,bold=True,color=RGBColor(0x15,0x65,0xC0))
for i,e in enumerate(['线速度:  v = (v_R + v_L) / 2','角速度:  ω = (v_R - v_L) / L','位姿变化:','  ẋ = v·cosθ    ẏ = v·sinθ    θ̇ = ω','轮角速度:','  ω_R = (2v + ωL) / 2r','  ω_L = (2v - ωL) / 2r']):
    txt(s6,16.0,2.8+i*0.45,16.0,0.4,e,sz=9,color=RGBColor(0x15,0x65,0xC0))
rrect(s6,15.5,7.5,17.0,3.5,RGBColor(0xFF,0xF3,0xE0),RGBColor(0xE6,0x51,0x00))
txt(s6,16.0,7.7,16.0,0.5,'底盘参数',sz=10,bold=True,color=RGBColor(0xE6,0x51,0x00))
for i,b in enumerate(['轮距 L = 0.52 m    轮半径 r = 0.09 m','标称速度 v = 0.3 m/s','最大角速度 ω = 1.5 rad/s','最小转弯半径 R = v/ω = 0.2 m']):
    txt(s6,16.0,8.3+i*0.45,16.0,0.4,b,sz=9,color=RGBColor(0xE6,0x51,0x00))
txt(s6,15.5,11.5,17.0,4.0,'差速驱动通过调节左右轮速度差实现转向：\n• 两侧轮速相等 → 直线行驶\n• 两侧轮速不等 → 向低速侧偏转\n• 原地转向：左右轮等速反向\n\n该模型是后续轨迹仿真、路径规划和PID控制的数学基础。',sz=9,color=GR)
footer(s6)
print("Slide 6: 运动学建模")

# ══════════════════════════════════════════
# Slide 7: 轨迹仿真
# ══════════════════════════════════════════
s7 = new_slide()
title_bar(s7,'05','典型轨迹仿真验证','5/11')
for path,name,l,t in [(os.path.join(fig,'直线.png'),'直线运动',1.0,2.0),(os.path.join(fig,'圆形.png'),'定曲率圆弧',9.0,2.0),(os.path.join(fig,'s型.png'),'S形运动',1.0,8.0),(os.path.join(fig,'8字.png'),'8字形运动',9.0,8.0)]:
    img(s7,path,l,t,w=7.5)
    txt(s7,l,t+5.8,7.5,0.4,name,sz=8,bold=True,color=DB,align=PP_ALIGN.CENTER)
rrect(s7,17.5,2.0,15.0,11.5,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s7,18.0,2.2,14.0,0.5,'轨迹仿真结果汇总',sz=11,bold=True,color=DB)
hdrs=['轨迹类型','路径/m','位移/m','ω_max','κ_max']
cx=[18.0,23.5,27.0,30.0,33.0]; cw=[5.0,3.0,2.5,2.5,2.5]
for i,h in enumerate(hdrs): txt(s7,cx[i],2.9,cw[i],0.4,h,sz=8,bold=True,color=DB)
rrect(s7,18.0,3.3,14.5,0.02,DB)
for ri,row in enumerate([('直线','9.0','9.0','0','0'),('定曲率圆','7.0','1.521','0.45','1.286'),('S形','8.0','5.269','0.60','1.500'),('8字形','7.6','5.817','0.70','1.842')]):
    y=3.5+ri*0.5
    rrect(s7,18.0,y,14.5,0.45,RGBColor(0xF5,0xF5,0xF5) if ri%2==0 else WH)
    for ci,val in enumerate(row): txt(s7,cx[ci],y+0.05,cw[ci],0.35,val,sz=8,color=GR)
txt(s7,18.0,5.8,14.0,0.4,'ω_max: rad/s    κ_max: m⁻¹    平均速度: 0.35~0.45 m/s',sz=7,color=LT)
img(s7,os.path.join(fig,'输入与轮速曲线.png'),17.5,6.5,w=7.5)
txt(s7,17.5,11.0,7.5,0.4,'控制输入与轮速曲线',sz=7,color=LT,align=PP_ALIGN.CENTER)
img(s7,os.path.join(fig,'曲率变化.png'),25.5,6.5,w=7.0)
txt(s7,25.5,11.0,7.0,0.4,'曲率对比',sz=7,color=LT,align=PP_ALIGN.CENTER)
txt(s7,1.0,14.5,32.0,2.5,'四类轨迹覆盖差速底盘的典型工况：直线体现基础直行能力，圆弧反映稳态转弯特性，S形和8字形呈现连续变曲率下的姿态调整能力。仿真验证了差速运动学模型的正确性，为路径规划和PID控制奠定基础。',sz=9,color=GR)
footer(s7)
print("Slide 7: 轨迹仿真")

# ══════════════════════════════════════════
# Slide 8: 路径规划与循迹
# ══════════════════════════════════════════
s8 = new_slide()
title_bar(s8,'06','路径规划与循迹仿真','6/11')
MX=1.0; CW=31.9
rrect(s8,MX,1.6,CW,7.2,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s8,MX+0.5,1.8,15,0.6,'Hybrid A* 仓储路径规划结果',sz=13,bold=True,color=DB)
img(s8,os.path.join(mlb,'hybrid_astar_result.png'),MX+0.3,2.6,w=13.0)
RX=15.0
txt(s8,RX,2.5,17.5,1.8,'算法原理：Hybrid A*将节点状态扩展为(x,y,θ)三维位姿，采用5种离散曲率的差速运动学基元扩展后继节点，代价函数f=g+h含转弯惩罚与倒车惩罚，使路径天然满足转弯半径约束。',sz=8.5,color=LT)
txt(s8,RX,4.4,8.0,0.5,'仿真参数',sz=10,bold=True,color=DB)
for i,p in enumerate(['地图 10m x 8m, 分辨率 0.2m','5 个货架障碍物','起点(1,1) -> 终点(9,7)','搜索迭代 3008 次']): txt(s8,RX,5.0+i*0.48,8.5,0.45,p,sz=8,color=GR)
for i,p in enumerate(['路径长度 11.04 m, 25 点','κ = {-1.11, -0.52, 0, 0.52, 1.11}','步长 0.5m, 轮距 L = 0.52m','3次均值平滑 + 0.05m等距插值']): txt(s8,RX+9.0,5.0+i*0.48,9.0,0.45,p,sz=8,color=GR)
txt(s8,RX,7.0,17.5,0.7,'规划结果：路径成功绕过所有货架，整体平滑无锯齿。',sz=8.5,bold=True,color=RGBColor(0x2E,0x7D,0x32))
txt(s8,MX,9.0,CW,0.4,'路径规划层 (Hybrid A*) 确定"走哪条路"  |  循迹执行层 (PID + 5路红外) 解决"如何沿路走"',sz=8,bold=True,color=DB,align=PP_ALIGN.CENTER)
LW=15.5; GAP=0.5; RW=CW-LW-GAP; IMG_H=5.5; IMG_T=10.2
rrect(s8,MX,9.5,LW,7.7,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s8,MX+0.5,9.6,10,0.5,'参考路径 vs 循迹轨迹',sz=11,bold=True,color=DB)
img(s8,os.path.join(mlb,'line_following_traj.png'),MX+0.3,IMG_T,h=IMG_H)
txt(s8,MX+0.5,16.0,LW-1.0,1.1,'蓝色实线为参考路径，红色虚线为循迹轨迹。两者高度重合。',sz=8,color=LT)
RX2=MX+LW+GAP
rrect(s8,RX2,9.5,RW,7.7,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s8,RX2+0.5,9.6,10,0.5,'循迹横向误差',sz=11,bold=True,color=DB)
img(s8,os.path.join(mlb,'line_following_error.png'),RX2+0.3,IMG_T,h=IMG_H)
txt(s8,RX2+0.5,16.0,RW-1.0,1.1,'起步误差较大，随后收敛至RMS=1.5cm，均在传感器覆盖裕量内。',sz=8,color=LT)
CY=17.5; CH=1.5
for i,(lab,val,bg,tc) in enumerate([('RMS横向误差','1.5 cm',RGBColor(0xE8,0xF5,0xE9),RGBColor(0x2E,0x7D,0x32)),('最大横向误差','6.8 cm',RGBColor(0xFF,0xF3,0xE0),RGBColor(0xE6,0x51,0x00)),('规划路径长度','11.04 m',RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0)),('循迹耗时','36.2 s',RGBColor(0xF3,0xE5,0xF5),RGBColor(0x7B,0x1F,0xA2))]):
    CW4=(CW-0.4*3)/4; cx=MX+0.2+i*(CW4+0.4)
    rrect(s8,cx,CY,CW4,CH,bg,tc)
    txt(s8,cx+0.2,CY+0.1,CW4-0.4,0.5,lab,sz=9,color=tc,align=PP_ALIGN.CENTER)
    txt(s8,cx+0.2,CY+0.6,CW4-0.4,0.7,val,sz=16,bold=True,color=tc,align=PP_ALIGN.CENTER)
footer(s8)
print("Slide 8: 路径规划")

# ══════════════════════════════════════════
# Slide 9: 运动基元与算法流程
# ══════════════════════════════════════════
s9 = new_slide()
title_bar(s9,'07','运动基元与算法流程','7/11')
rrect(s9,1.0,2.0,15.5,7.0,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s9,1.5,2.2,14.5,0.5,'差速运动基元 (5种曲率)',sz=11,bold=True,color=DB)
for i,(k,rn,name,color) in enumerate([('κ=-1.11','R=0.90m','大左转',RGBColor(0xE6,0x51,0x00)),('κ=-0.52','R=1.92m','小左转',RGBColor(0xF3,0x9C,0x12)),('κ=0','R=∞','直行',RGBColor(0x27,0xAE,0x60)),('κ=+0.52','R=1.92m','小右转',RGBColor(0x29,0x80,0xB9)),('κ=+1.11','R=0.90m','大右转',RGBColor(0x8E,0x44,0xAD))]):
    y=3.0+i*0.7
    rrect(s9,1.5,y,14.0,0.6,RGBColor(0xFA,0xFA,0xFA),color)
    txt(s9,2.0,y+0.1,4.0,0.4,name,sz=9,bold=True,color=color)
    txt(s9,6.0,y+0.1,4.5,0.4,k,sz=8,color=GR)
    txt(s9,10.5,y+0.1,4.5,0.4,rn,sz=8,color=GR)
txt(s9,1.5,6.8,14.5,1.5,'前进/后退 × 5种曲率 = 10种运动基元\n步长 d = 0.5m\n节点状态: (x, y, θ) 三维位姿',sz=8,color=LT)
rrect(s9,17.0,2.0,15.5,7.0,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s9,17.5,2.2,14.5,0.5,'Hybrid A* 算法流程',sz=11,bold=True,color=DB)
for i,(name,desc,bg) in enumerate([('栅格地图','10×8m, 0.2m分辨率',RGBColor(0xE8,0xEA,0xF6)),('初始化','Open/Close列表, 启发函数h₀',RGBColor(0xE8,0xEA,0xF6)),('节点扩展','10种运动基元生成后继',RGBColor(0xE8,0xEA,0xF6)),('碰撞检测','插值采样11个点',RGBColor(0xE8,0xEA,0xF6)),('代价评估','f=g+h, 含转弯/倒车惩罚',RGBColor(0xE8,0xEA,0xF6)),('路径回溯','从终点回溯父节点链',RGBColor(0xE8,0xF5,0xE9)),('均值平滑','3次加权平均滤波',RGBColor(0xE8,0xF5,0xE9)),('等距插值','0.05m间距, 224个路径点',RGBColor(0xE8,0xF5,0xE9))]):
    y=2.9+i*0.55
    rrect(s9,17.5,y,14.5,0.48,bg)
    txt(s9,18.0,y+0.05,4.0,0.38,f'{i+1}. {name}',sz=8,bold=True,color=DB)
    txt(s9,22.5,y+0.05,9.0,0.38,desc,sz=8,color=GR)
rrect(s9,1.0,9.5,31.5,3.0,RGBColor(0xFF,0xF8,0xE1),RGBColor(0xF9,0xA8,0x25))
txt(s9,1.5,9.7,30.5,0.5,'代价函数 f = g + h',sz=11,bold=True,color=RGBColor(0xF5,0x7F,0x17))
txt(s9,1.5,10.3,10.0,2.0,'g: 实际代价\n  • 基础: 步长0.5m\n  • 倒车惩罚: +50%\n  • 转弯惩罚: +0.1',sz=9,color=RGBColor(0xF5,0x7F,0x17))
txt(s9,12.0,10.3,10.0,2.0,'h: 启发值\n  • 欧氏距离\n  • 保证可采纳性',sz=9,color=RGBColor(0xF5,0x7F,0x17))
txt(s9,23.0,10.3,9.0,2.0,'关键特性:\n• 节点含航向θ\n• 满足差速约束\n• 路径可直接执行',sz=9,bold=True,color=RGBColor(0xF5,0x7F,0x17))
img(s9,os.path.join(fig,'sensor_layout.png'),1.0,13.0,w=6.0)
txt(s9,1.0,17.0,6.0,0.4,'5路红外传感器布局',sz=7,color=LT,align=PP_ALIGN.CENTER)
txt(s9,7.5,13.2,25.0,4.0,'循迹传感器阵列：\n• 5路TCRT5000红外，间距40mm，总覆盖160mm\n• 安装在车体前端、驱动轮轴线前方80mm\n• 权重{-4,-2,0,+2,+4}计算加权偏移量\n• 丢线时外推1.5倍上次偏移量\n• PID: Kp=3.0, Ki=0.08, Kd=1.0, 限幅±1.5rad/s',sz=9,color=GR)
footer(s9)
print("Slide 9: 运动基元")

# ══════════════════════════════════════════
# Slide 10: 双环PID控制
# ══════════════════════════════════════════
s10 = new_slide()
title_bar(s10,'08','双环PID运动控制','8/11')
img(s10,os.path.join(mlb,'cascaded_pid_diagram.png'),1.0,2.0,w=11.5)
txt(s10,1.0,10.5,11.5,0.4,'双环PID级联控制结构',sz=8,color=LT,align=PP_ALIGN.CENTER)
rrect(s10,13.0,2.0,19.5,4.0,RGBColor(0xE3,0xF2,0xFD),RGBColor(0x15,0x65,0xC0))
txt(s10,13.5,2.2,18.5,0.5,'双环PID控制器参数',sz=10,bold=True,color=RGBColor(0x15,0x65,0xC0))
px=[13.5,20.0,23.5,27.0,30.0]; pw=[6.0,3.0,3.0,2.5,3.0]
for i,h in enumerate(['控制器','Kp','Ki','Kd','说明']): txt(s10,px[i],2.8,pw[i],0.4,h,sz=8,bold=True,color=DB)
for ri,row in enumerate([('位姿环(航向PID)','3.0','0.2','0.4','外环,输出v,ω'),('轮速环(PI)','17.0','103.0','--','内环,输出PWM')]):
    y=3.3+ri*0.5
    rrect(s10,13.5,y,18.5,0.45,RGBColor(0xF5,0xF5,0xF5) if ri%2==0 else WH)
    for ci,val in enumerate(row): txt(s10,px[ci],y+0.05,pw[ci],0.35,val,sz=8,color=GR)
txt(s10,13.5,4.6,18.5,1.2,'内环响应频率约为外环3~5倍，保证内环优先收敛。\n轮速环采用PI而非PID：微分项对编码器噪声敏感。',sz=8,color=LT)
for path,name,l,t in [(os.path.join(mlb,'pid_step_response.png'),'阶跃响应',13.0,6.5),(os.path.join(mlb,'pid_point_stabilization.png'),'定点镇定',22.0,6.5),(os.path.join(mlb,'pid_line_tracking.png'),'直线跟踪',13.0,11.5),(os.path.join(mlb,'pid_circle_tracking.png'),'圆形跟踪',22.0,11.5)]:
    img(s10,path,l,t,w=8.5)
    txt(s10,l,t+6.0,8.5,0.4,name,sz=8,bold=True,color=DB,align=PP_ALIGN.CENTER)
rrect(s10,1.0,13.5,11.5,3.0,WH,RGBColor(0xE0,0xE0,0xE0))
txt(s10,1.5,13.7,10.5,0.5,'仿真结果汇总',sz=10,bold=True,color=DB)
for i,(sc,res) in enumerate([('阶跃响应','稳态误差≈0, 调节时间0.8s'),('定点镇定','位置误差(0.022,-0.007)m, 航向-0.89°'),('直线跟踪','角速度RMSE 0.0023 rad/s'),('圆形跟踪','线速度RMSE 0.0086m/s, 圆度RMSE 0.020m')]):
    y=14.3+i*0.5
    txt(s10,1.5,y,3.0,0.4,sc,sz=8,bold=True,color=DB)
    txt(s10,4.5,y,7.5,0.4,res,sz=8,color=GR)
footer(s10)
print("Slide 10: 双环PID")

# ══════════════════════════════════════════
# Slide 11: 致谢
# ══════════════════════════════════════════
s11 = new_slide()
rrect(s11,0,0,33.9,19.1,RGBColor(0x1A,0x23,0x7E))
rrect(s11,0,0,33.9,3.5,RGBColor(0x28,0x35,0x93))
txt(s11,2,6,30,1.5,'感谢各位老师指导',sz=32,bold=True,color=WH,align=PP_ALIGN.CENTER)
rrect(s11,8,8,18,0.03,RGBColor(0x3F,0x51,0xB5))
txt(s11,2,9.5,30,0.6,'答辩人：王梓煜    指导教师：曹树春',sz=14,color=RGBColor(0xBB,0xDE,0xFB),align=PP_ALIGN.CENTER)
txt(s11,2,10.5,30,0.6,'齐鲁工业大学  机械工程学院',sz=14,color=RGBColor(0xBB,0xDE,0xFB),align=PP_ALIGN.CENTER)
print("Slide 11: 致谢")

# ══════════════════════════════════════════
# 保存
# ══════════════════════════════════════════
prs.save(new_pptx)
print(f"\n已保存: {new_pptx}")
print(f"共 {len(prs.slides)} 页幻灯片")
