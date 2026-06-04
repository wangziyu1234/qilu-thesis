"""答辩PPT第7页：路径规划与循迹 — 无重叠布局"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

pptx_path = 'D:/bylw/code/答辩PPT/答辩PPT_图文并茂.pptx'
fig_dir = 'D:/bylw/code/fangzhen/matlab'

prs = Presentation(pptx_path)
slide = prs.slides[6]

DB = RGBColor(0x28,0x35,0x93)
WH = RGBColor(0xFF,0xFF,0xFF)
GR = RGBColor(0x66,0x66,0x66)
LT = RGBColor(0x55,0x55,0x55)

def rrect(l,t,w,h,fill,line=None):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Cm(l),Cm(t),Cm(w),Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line: s.line.color.rgb=line; s.line.width=Pt(1)
    else: s.line.fill.background()
    return s

def txt(l,t,w,h,text,sz=12,bold=False,color=GR,align=PP_ALIGN.LEFT):
    tb=slide.shapes.add_textbox(Cm(l),Cm(t),Cm(w),Cm(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text
    p.font.size=Pt(sz); p.font.bold=bold; p.font.color.rgb=color; p.alignment=align

def img(path,l,t,w=None,h=None):
    if w and h: return slide.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w),height=Cm(h))
    elif w: return slide.shapes.add_picture(path,Cm(l),Cm(t),width=Cm(w))
    elif h: return slide.shapes.add_picture(path,Cm(l),Cm(t),height=Cm(h))
    else: return slide.shapes.add_picture(path,Cm(l),Cm(t))

# ===== 清除旧内容(保留标题栏) =====
keep = {'Shape 0','Shape 1','Shape 2','Text 3','Text 4','Text 6'}
for s in [sh for sh in slide.shapes if sh.name not in keep]:
    s._element.getparent().remove(s._element)

MX = 1.0; CW = 31.9

# ══════════════════════════════════════════════
#  上区：Y=1.6 ~ 8.8 (7.2cm)
#  左: 地图(宽14cm)  右: 算法说明+参数(宽17cm)
# ══════════════════════════════════════════════
rrect(MX, 1.6, CW, 7.2, WH, RGBColor(0xE0,0xE0,0xE0))
txt(MX+0.5, 1.8, 15, 0.6, 'Hybrid A* 仓储路径规划结果', sz=13, bold=True, color=DB)

# 地图(左侧, 固定宽13cm, 高度自动)
img(os.path.join(fig_dir,'hybrid_astar_result.png'), MX+0.3, 2.6, w=13.0)

# 右侧区域 X=15.0 ~ 32.9
RX = 15.0

# 算法说明(右上)
txt(RX, 2.5, 17.5, 1.8,
    '算法原理：Hybrid A*将节点状态扩展为(x,y,θ)三维位姿，'
    '采用5种离散曲率的差速运动学基元扩展后继节点，'
    '代价函数f=g+h含转弯惩罚与倒车惩罚，'
    '使路径天然满足转弯半径约束，无需后处理即可执行。',
    sz=8.5, color=LT)

# 参数(右中)
txt(RX, 4.4, 8.0, 0.5, '仿真参数', sz=10, bold=True, color=DB)
params_l = [
    '地图 10m x 8m, 分辨率 0.2m',
    '5 个货架障碍物',
    '起点(1,1) -> 终点(9,7)',
    '搜索迭代 3008 次',
]
params_r = [
    '路径长度 11.04 m, 25 点',
    'κ = {-1.11, -0.52, 0, 0.52, 1.11}',
    '步长 0.5m, 轮距 L = 0.52m',
    '3次均值平滑 + 0.05m等距插值',
]
for i,p in enumerate(params_l):
    txt(RX, 5.0+i*0.48, 8.5, 0.45, p, sz=8, color=GR)
for i,p in enumerate(params_r):
    txt(RX+9.0, 5.0+i*0.48, 9.0, 0.45, p, sz=8, color=GR)

# 规划结论(右下)
txt(RX, 7.0, 17.5, 0.7,
    '规划结果：路径成功绕过所有货架，在障碍物间隙处合理转弯，整体平滑无锯齿。',
    sz=8.5, bold=True, color=RGBColor(0x2E,0x7D,0x32))

# ══════════════════════════════════════════════
#  分隔条 Y=9.0 ~ 9.4
# ══════════════════════════════════════════════
txt(MX, 9.0, CW, 0.4,
    '路径规划层 (Hybrid A*) 确定"走哪条路"    |    循迹执行层 (PID + 5路红外) 解决"如何沿路走"',
    sz=8, bold=True, color=DB, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
#  下区：Y=9.5 ~ 17.2 (7.7cm)
#  图片固定高度5.5cm, 说明文字在图片下方
# ══════════════════════════════════════════════
IMG_H = 5.5   # 图片固定高度
IMG_T = 10.2  # 图片起始Y
IMG_B = IMG_T + IMG_H  # = 15.7

LW = 15.5; GAP = 0.5; RW = CW - LW - GAP

# --- 左卡片：循迹轨迹 ---
rrect(MX, 9.5, LW, 7.7, WH, RGBColor(0xE0,0xE0,0xE0))
txt(MX+0.5, 9.6, 10, 0.5, '参考路径 vs 循迹轨迹', sz=11, bold=True, color=DB)
img(os.path.join(fig_dir,'line_following_traj.png'), MX+0.3, IMG_T, h=IMG_H)
# 说明(图片下方, Y=16.0)
txt(MX+0.5, 16.0, LW-1.0, 1.1,
    '蓝色实线为参考路径，红色虚线为循迹轨迹。'
    '两者高度重合，PID方案能稳定跟踪规划路径，转弯区段偏移小。',
    sz=8, color=LT)

# --- 右卡片：误差曲线 ---
RX2 = MX + LW + GAP
rrect(RX2, 9.5, RW, 7.7, WH, RGBColor(0xE0,0xE0,0xE0))
txt(RX2+0.5, 9.6, 10, 0.5, '循迹横向误差', sz=11, bold=True, color=DB)
img(os.path.join(fig_dir,'line_following_error.png'), RX2+0.3, IMG_T, h=IMG_H)
# 说明(图片下方, Y=16.0)
txt(RX2+0.5, 16.0, RW-1.0, 1.1,
    '起步误差较大(初始位姿偏差)，随后快速收敛至RMS=0.96cm。'
    '弯曲区段误差略增，均在传感器覆盖宽度(+80mm)裕量内。',
    sz=8, color=LT)

# ══════════════════════════════════════════════
#  底部指标卡片 Y=17.5 ~ 19.0
# ══════════════════════════════════════════════
CY = 17.5; CH = 1.5
cards = [
    ('RMS 横向误差', '0.96 cm', RGBColor(0xE8,0xF5,0xE9), RGBColor(0x2E,0x7D,0x32)),
    ('最大横向误差', '3.01 cm', RGBColor(0xFF,0xF3,0xE0), RGBColor(0xE6,0x51,0x00)),
    ('规划路径长度', '11.04 m', RGBColor(0xE3,0xF2,0xFD), RGBColor(0x15,0x65,0xC0)),
    ('循迹耗时',     '36.2 s',  RGBColor(0xF3,0xE5,0xF5), RGBColor(0x7B,0x1F,0xA2)),
]
CW4 = (CW - 0.4*3) / 4
for i,(lab,val,bg,tc) in enumerate(cards):
    cx = MX + 0.2 + i*(CW4+0.4)
    rrect(cx, CY, CW4, CH, bg, tc)
    txt(cx+0.2, CY+0.1, CW4-0.4, 0.5, lab, sz=9, color=tc, align=PP_ALIGN.CENTER)
    txt(cx+0.2, CY+0.6, CW4-0.4, 0.7, val, sz=16, bold=True, color=tc, align=PP_ALIGN.CENTER)

tmp_path = pptx_path + '.tmp'
prs.save(tmp_path)
print(f"已保存临时文件: {tmp_path}")

import shutil
shutil.move(tmp_path, pptx_path)
print(f"已覆盖: {pptx_path}")
