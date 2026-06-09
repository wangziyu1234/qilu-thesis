"""答辩PPT图文调整：插入仿真图片，合理布局"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

SRC = r"D:\bylw\code\答辩PPT\答辩PPT_完整版.pptx"
DST = r"D:\bylw\code\答辩PPT\答辩PPT_完整版_图文调整.pptx"
FIG_DIR = r"D:\code\matlab\figures"

prs = Presentation(SRC)

# 幻灯片尺寸 (标准16:9 ≈ 33.87cm × 19.05cm)
SW = prs.slide_width   # 12192000 EMU
SH = prs.slide_height  # 6858000 EMU

def add_image(slide, filename, left_cm, top_cm, width_cm):
    """添加图片到指定位置，自动计算高度保持比例"""
    path = os.path.join(FIG_DIR, filename)
    if not os.path.exists(path):
        print(f"  ⚠ 图片不存在: {path}")
        return None
    left = Cm(left_cm)
    top = Cm(top_cm)
    width = Cm(width_cm)
    pic = slide.shapes.add_picture(path, left, top, width=width)
    print(f"  ✓ {filename} → ({left_cm},{top_cm}) w={width_cm}cm")
    return pic

# ===== 各幻灯片图文调整 =====

# Slide 7 (index 6): 机械结构——剪叉机构 → 加机构简图
slide7 = prs.slides[6]
add_image(slide7, "jiagou_low.png", 17.5, 5.5, 14.5)
add_image(slide7, "jiagou_high.png", 17.5, 12.0, 14.5)

# Slide 8 (index 7): 强度校核 → 加FEA结果图
slide8 = prs.slides[7]
add_image(slide8, "fea_results_50kg.png", 17.0, 5.5, 15.0)

# Slide 10 (index 9): 传感器与电路 → 加传感器布局图
slide10 = prs.slides[9]
add_image(slide10, "sensor_layout.png", 18.0, 12.0, 10.0)

# Slide 11 (index 10): 运动学建模与仿真 → 加轨迹图
slide11 = prs.slides[10]
add_image(slide11, "直线.png", 17.5, 5.5, 7.0)
add_image(slide11, "圆形.png", 25.0, 5.5, 7.0)
add_image(slide11, "s型.png", 17.5, 12.0, 7.0)
add_image(slide11, "8字.png", 25.0, 12.0, 7.0)

# Slide 12 (index 11): 路径规划与循迹
slide12 = prs.slides[11]
add_image(slide12, "hybrid_astar_result.png", 17.5, 5.0, 15.0)
add_image(slide12, "line_following_traj.png", 17.5, 13.0, 10.0)

# Slide 13 (index 12): 运动控制算法设计 → 加PID框图
slide13 = prs.slides[12]
add_image(slide13, "cascaded_pid_diagram.png", 17.0, 5.0, 15.5)

# Slide 14 (index 13): 运动控制仿真验证 → 加PID仿真结果
slide14 = prs.slides[13]
add_image(slide14, "仿真结果_场景2_定点镇定.png", 17.5, 5.0, 15.0)
add_image(slide14, "仿真结果_场景4_圆形跟踪.png", 17.5, 13.0, 15.0)

# Slide 6 (index 5): 总体方案-控制与规划 → 加路径规划示意图
slide6 = prs.slides[5]
add_image(slide6, "cascaded_pid_diagram.png", 17.0, 5.5, 14.0)

# Slide 5 (index 4): 总体方案-驱动与升降 → 加驱动示意
slide5 = prs.slides[4]
add_image(slide5, "jiagou_low.png", 17.5, 5.5, 13.0)

prs.save(DST)
print(f"\n完成！已保存至: {DST}")
